import streamlit as st
import cohere
import PyPDF2
import pandas as pd
from docx import Document
from pptx import Presentation
import pytesseract
from PIL import Image
import os
from config import config   

COHERE_API_KEY = config["COHERE_API_KEY"]   
co = cohere.Client(COHERE_API_KEY)

def read_txt(file):
    return file.read().decode("utf-8")

def read_pdf(file):
    reader = PyPDF2.PdfReader(file)
    return "\n".join([page.extract_text() for page in reader.pages])

def read_excel(file):
    data = pd.read_excel(file)
    return data.to_string()

def read_csv(file):
    data = pd.read_csv(file)
    return data.to_string()

def read_docx(file):
    doc = Document(file)
    return "\n".join([para.text for para in doc.paragraphs])

def read_pptx(file):
    presentation = Presentation(file)
    return "\n".join([shape.text for slide in presentation.slides for shape in slide.shapes if shape.has_text_frame])

def read_image(file):
    image = Image.open(file)
    return pytesseract.image_to_string(image)

def read_file(uploaded_file):
    ext = uploaded_file.name.split(".")[-1]
    if ext in ["txt"]:
        return read_txt(uploaded_file)
    elif ext in ["pdf"]:
        return read_pdf(uploaded_file)
    elif ext in ["xls", "xlsx"]:
        return read_excel(uploaded_file)
    elif ext in ["csv"]:
        return read_csv(uploaded_file)
    elif ext in ["docx"]:
        return read_docx(uploaded_file)
    elif ext in ["pptx"]:
        return read_pptx(uploaded_file)
    elif ext in ["jpg", "jpeg", "png", "tiff"]:
        return read_image(uploaded_file)
    return "Unsupported file type."

def chatbot_response(query, document_text):
    prompt = f"""
    Document content:\n\n{document_text}\n\n
    Question: {query}
    Answer concisely and directly.
    """
    response = co.generate(
        model="command-xlarge-nightly",
        prompt=prompt,
        max_tokens=5000,
        temperature=0.7
    )
    return response.generations[0].text.strip()

st.title("Document-Based Chatbot")
st.write("Upload a document, ask a question, and get answers!")

uploaded_file = st.file_uploader("Upload your document", type=['txt', 'pdf', 'docx', 'xlsx', 'csv', 'pptx', 'jpg', 'jpeg', 'png', 'tiff'])

document_text = ""
if uploaded_file:
    document_text = read_file(uploaded_file)
    st.success("Document processed successfully!")

user_query = st.text_input("Enter your question about the document:")
if user_query and document_text:
    response = chatbot_response(user_query, document_text)
    st.write("### Chatbot Response:")
    st.write(response)
