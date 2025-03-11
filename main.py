import streamlit as st
import cohere
import PyPDF2
import pandas as pd
from docx import Document
from pptx import Presentation
import pytesseract
from PIL import Image
import os
from pyngrok import ngrok
import threading
from config import config  # Import configuration

# Initialize Cohere API
COHERE_API_KEY = config["COHERE_API_KEY"]  # Fetch Cohere API key from config
co = cohere.Client(COHERE_API_KEY)

# File parsers (no changes here)
def read_txt(file_path):
    with open(file_path, 'r', encoding='utf-8') as file:
        return file.read()

def read_pdf(file_path):
    content = ""
    with open(file_path, 'rb') as file:
        reader = PyPDF2.PdfReader(file)
        for page in reader.pages:
            content += page.extract_text()
    return content

def read_excel(file_path):
    data = pd.read_excel(file_path)
    return data.to_string()

def read_csv(file_path):
    data = pd.read_csv(file_path)
    return data.to_string()

def read_docx(file_path):
    doc = Document(file_path)
    content = ""
    for paragraph in doc.paragraphs:
        content += paragraph.text + "\n"
    return content

def read_pptx(file_path):
    presentation = Presentation(file_path)
    content = ""
    for slide in presentation.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                content += shape.text + "\n"
    return content

def read_image(file_path):
    image = Image.open(file_path)
    return pytesseract.image_to_string(image)

def read_file(file_path):
    if file_path.endswith('.txt'):
        return read_txt(file_path)
    elif file_path.endswith('.pdf'):
        return read_pdf(file_path)
    elif file_path.endswith(('.xls', '.xlsx')):
        return read_excel(file_path)
    elif file_path.endswith('.csv'):
        return read_csv(file_path)
    elif file_path.endswith('.docx'):
        return read_docx(file_path)
    elif file_path.endswith('.pptx'):
        return read_pptx(file_path)
    elif file_path.endswith(('.jpg', '.jpeg', '.png', '.tiff')):
        return read_image(file_path)
    else:
        return "Unsupported file type."

# Chatbot response function
def chatbot_response(query, document_text):
    prompt = f"""
    Document content:\n\n{document_text}\n\n
    Question: {query}
    Answer concisely and directly.
    """
    response = co.generate(
        model="command-xlarge-nightly",
        prompt=prompt,
        max_tokens=150,
        temperature=0.7
    )
    return response.generations[0].text.strip()

# Streamlit App
def run_streamlit():
    st.title("Document-Based Chatbot")
    st.write("Upload a document, ask a question, and get answers!")

    # File upload
    uploaded_file = st.file_uploader("Upload your document", type=['txt', 'pdf', 'docx', 'xlsx', 'csv', 'pptx', 'jpg', 'jpeg', 'png', 'tiff'])

    if uploaded_file:
        # Save the file temporarily
        file_path = os.path.join("temp", uploaded_file.name)
        os.makedirs("temp", exist_ok=True)
        with open(file_path, "wb") as f:
            f.write(uploaded_file.getbuffer())

        # Process the file
        document_text = read_file(file_path)
        st.success("Document processed successfully!")

        # User query
        user_query = st.text_input("Enter your question about the document:")
        if user_query:
            response = chatbot_response(user_query, document_text)
            st.write("### Chatbot Response:")
            st.write(response)

# Ngrok Integration
def start_ngrok():
    ngrok.set_auth_token(config["NGROK_AUTH_TOKEN"])  # Fetch Ngrok auth token from config
    public_url = ngrok.connect(8501)
    print(f"Streamlit app is live at: {public_url}")

if __name__ == "__main__":
    # Start Ngrok in a separate thread
    thread = threading.Thread(target=start_ngrok)
    thread.start()

    # Run the Streamlit app
    run_streamlit()
